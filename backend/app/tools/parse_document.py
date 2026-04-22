"""
parse_document 工具 — 解析上传的文档文件。
支持格式: PDF / DOCX / PPTX / XLSX / TXT / MD / CSV。
解析后返回结构化文本内容 + 元数据，可选通过 memory_service 索引。
"""
import csv
import io
import logging
import os
from pathlib import Path, PurePosixPath
from typing import Any

logger = logging.getLogger(__name__)

# ──────────────── Tool 定义（OpenAI function-calling 格式）────────────────
TOOL_DEFINITION: dict[str, Any] = {
    "type": "function",
    "function": {
        "name": "parse_document",
        "description": (
            "解析上传的文档文件，提取文本内容。"
            "支持 PDF、DOCX、PPTX、XLSX、TXT、MD、CSV 格式。"
            "返回文档的文本内容和分块，供后续分析或向量索引使用。"
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "asset_id": {
                    "type": "string",
                    "description": "上传文件的 Asset ID",
                },
                "file_path": {
                    "type": "string",
                    "description": "文件在服务器上的路径（由上传 API 返回的 file_url 对应路径）",
                },
                "max_chars": {
                    "type": "integer",
                    "description": "返回的最大字符数（默认 50000），超出时截断并提示",
                    "default": 50000,
                },
                "index_chunks": {
                    "type": "boolean",
                    "description": "是否将解析结果索引到向量数据库（默认 false）",
                    "default": False,
                },
            },
            "required": ["asset_id", "file_path"],
        },
    },
}

# 分块大小（用于向量索引）
CHUNK_SIZE = 800       # 每块约 800 字符
CHUNK_OVERLAP = 100    # 块间重叠字符数


# ──────────────────── 解析函数 ────────────────────

def _resolve_path(file_path: str) -> str:
    """
    将 file_url 格式（/static/uploads/...）转为实际磁盘路径。
    /static → data/
    """
    backend_root = Path(__file__).resolve().parents[2]

    if file_path.startswith("/static/"):
        relative_path = PurePosixPath(file_path.replace("/static/", "data/", 1))
        return str((backend_root / Path(relative_path)).resolve())

    raw_path = Path(file_path)
    if raw_path.is_absolute():
        return str(raw_path)

    return str((backend_root / raw_path).resolve())


def _parse_pdf(path: str) -> dict[str, Any]:
    """解析 PDF 文件（使用 PyMuPDF）。"""
    import fitz  # PyMuPDF

    doc = fitz.open(path)
    pages: list[dict[str, Any]] = []
    full_text = []

    for i, page in enumerate(doc):
        text = page.get_text("text")
        pages.append({
            "page": i + 1,
            "text": text,
            "char_count": len(text),
        })
        full_text.append(text)

    doc.close()

    return {
        "format": "pdf",
        "page_count": len(pages),
        "pages": pages,
        "full_text": "\n\n".join(full_text),
    }


def _parse_docx(path: str) -> dict[str, Any]:
    """解析 DOCX 文件（使用 python-docx）。"""
    from docx import Document

    doc = Document(path)
    paragraphs: list[str] = []

    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)

    # 提取表格内容
    tables_text: list[str] = []
    for table in doc.tables:
        rows_data = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows_data.append(" | ".join(cells))
        tables_text.append("\n".join(rows_data))

    full_text = "\n\n".join(paragraphs)
    if tables_text:
        full_text += "\n\n[表格内容]\n" + "\n\n".join(tables_text)

    return {
        "format": "docx",
        "paragraph_count": len(paragraphs),
        "table_count": len(doc.tables),
        "full_text": full_text,
    }


def _parse_pptx(path: str) -> dict[str, Any]:
    """解析 PPTX 文件（使用 python-pptx）。"""
    from pptx import Presentation

    prs = Presentation(path)
    slides_data: list[dict[str, Any]] = []
    full_text = []

    for i, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)

        slide_text = "\n".join(texts)
        slides_data.append({
            "slide": i + 1,
            "text": slide_text,
        })
        full_text.append(f"[幻灯片 {i + 1}]\n{slide_text}")

    return {
        "format": "pptx",
        "slide_count": len(slides_data),
        "slides": slides_data,
        "full_text": "\n\n".join(full_text),
    }


def _parse_xlsx(path: str) -> dict[str, Any]:
    """解析 XLSX 文件（使用 openpyxl）。"""
    from openpyxl import load_workbook

    wb = load_workbook(path, data_only=True, read_only=True)
    sheets_data: list[dict[str, Any]] = []
    full_text = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):  # 跳过全空行
                rows.append(" | ".join(cells))

        sheet_text = "\n".join(rows)
        sheets_data.append({
            "sheet": sheet_name,
            "row_count": len(rows),
            "text": sheet_text,
        })
        full_text.append(f"[工作表: {sheet_name}]\n{sheet_text}")

    wb.close()

    return {
        "format": "xlsx",
        "sheet_count": len(sheets_data),
        "sheets": sheets_data,
        "full_text": "\n\n".join(full_text),
    }


def _parse_csv(path: str) -> dict[str, Any]:
    """解析 CSV 文件。"""
    rows: list[str] = []
    with open(path, 'r', encoding='utf-8', errors='replace') as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(" | ".join(row))

    return {
        "format": "csv",
        "row_count": len(rows),
        "full_text": "\n".join(rows),
    }


def _parse_text(path: str) -> dict[str, Any]:
    """解析纯文本文件（TXT / MD / 代码文件）。"""
    with open(path, 'r', encoding='utf-8', errors='replace') as f:
        content = f.read()

    ext = Path(path).suffix.lower()
    fmt = {".md": "markdown", ".txt": "text"}.get(ext, "text")

    return {
        "format": fmt,
        "line_count": content.count('\n') + 1,
        "full_text": content,
    }


# 扩展名 → 解析器映射
_PARSERS: dict[str, callable] = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
    ".xlsx": _parse_xlsx,
    ".csv": _parse_csv,
    ".txt": _parse_text,
    ".md": _parse_text,
}


def _split_chunks(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[dict]:
    """将文本分块（用于向量索引）。"""
    chunks = []
    start = 0
    idx = 0
    while start < len(text):
        end = start + chunk_size
        chunk_text = text[start:end]
        chunks.append({
            "index": idx,
            "content": chunk_text,
            "metadata": {"start": start, "end": end},
        })
        idx += 1
        start = end - overlap
    return chunks


# ──────────────── Tool 执行入口 ────────────────

async def execute(params: dict[str, Any]) -> dict[str, Any]:
    """
    执行文档解析。

    Args:
        params: {
            "asset_id": str,
            "file_path": str,
            "max_chars": int (可选, 默认 50000),
            "index_chunks": bool (可选, 默认 False),
        }

    Returns:
        {
            "asset_id": str,
            "format": str,
            "content": str,      # 解析出的文本（可能截断）
            "truncated": bool,
            "metadata": dict,    # 页数/行数等
            "chunks_indexed": int | None,
        }
    """
    asset_id = params.get("asset_id", "")
    file_path = params.get("file_path", "")
    max_chars = params.get("max_chars", 50000)
    index_chunks = params.get("index_chunks", False)

    if not asset_id or not file_path:
        return {"error": "缺少 asset_id 或 file_path 参数"}

    # 解析实际文件路径
    actual_path = _resolve_path(file_path)

    if not os.path.isfile(actual_path):
        return {"error": f"文件不存在: {file_path}"}

    # 确定解析器
    ext = Path(actual_path).suffix.lower()
    parser = _PARSERS.get(ext)

    if parser is None:
        # 尝试作为纯文本解析
        parser = _parse_text

    try:
        result = parser(actual_path)
    except Exception as e:
        logger.exception(f"[parse_document] 解析失败: {actual_path}")
        return {"error": f"文档解析失败: {str(e)}"}

    # 提取文本并检查截断
    full_text = result.pop("full_text", "")
    truncated = len(full_text) > max_chars
    content = full_text[:max_chars]

    if truncated:
        content += f"\n\n... [内容已截断，总共 {len(full_text)} 字符，显示前 {max_chars} 字符]"

    # 可选: 索引到向量数据库
    chunks_indexed = None
    if index_chunks and full_text.strip():
        try:
            from app.models.database import async_session
            from app.services.memory_service import index_document_chunks

            chunks = _split_chunks(full_text)
            async with async_session() as session:
                chunks_indexed = await index_document_chunks(
                    session=session,
                    asset_id=asset_id,
                    chunks=chunks,
                )
            logger.info(f"[parse_document] 向量索引完成: {chunks_indexed} 块")
        except Exception as e:
            logger.warning(f"[parse_document] 向量索引失败（不影响解析结果）: {e}")

    return {
        "asset_id": asset_id,
        "format": result.get("format", ext.lstrip(".")),
        "content": content,
        "truncated": truncated,
        "char_count": len(full_text),
        "metadata": {k: v for k, v in result.items() if k != "full_text"},
        "chunks_indexed": chunks_indexed,
    }
