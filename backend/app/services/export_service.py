"""
导出服务 — 支持多种格式导出演示文稿。
格式:
    - HTML: 完整 reveal.js 独立网页
    - PDF:  通过 Playwright 渲染后导出 PDF
    - PPTX 保真: Playwright 截图 → python-pptx 图片幻灯片（视觉一致）
    - PPTX 可编辑: 解析 HTML → python-pptx 文本/形状幻灯片（内容可编辑）
    - PPTX Native: DeckSpec -> Node/PptxGenJS renderer（推荐）
"""
import logging
import os
import re
import uuid
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from sqlalchemy.ext.asyncio import AsyncSession

from app.services.package_runtime import (
    DEFAULT_RUNTIME_USER_ID,
    OFFICIAL_NATIVE_ORCHESTRATOR,
    invoke_pptx_workflow_package,
)
from app.services.plugin_registry import record_artifact_variant
from app.services.theme_manager import get_theme

logger = logging.getLogger(__name__)

BACKEND_ROOT = Path(__file__).resolve().parents[2]

# 导出文件目录
EXPORT_DIR = BACKEND_ROOT / "data" / "exports"
EXPORT_DIR.mkdir(parents=True, exist_ok=True)


async def export_html(
    full_html: str,
    title: str,
) -> str:
    """
    导出 HTML 文件。

    Args:
        full_html: 完整的 reveal.js HTML 字符串
        title: 文件名前缀

    Returns:
        导出文件的相对路径
    """
    safe_title = re.sub(r'[^\w\u4e00-\u9fff-]', '_', title)[:50]
    filename = f"{safe_title}_{uuid.uuid4().hex[:8]}.html"
    filepath = EXPORT_DIR / filename

    filepath.write_text(full_html, encoding="utf-8")
    logger.info(f"[Export] HTML 导出完成: {filepath}")

    return f"exports/{filename}"


async def export_pdf(
    full_html: str,
    title: str,
) -> str:
    """
    通过 Playwright 渲染页面并导出 PDF。
    支持两种 HTML 结构：
    1. reveal.js: .reveal .slides > section
    2. webdeck: .deck-slide + #slides-container

    通过 JavaScript 动态修改 CSS，让所有幻灯片垂直排列后再生成 PDF。

    Args:
        full_html: 完整的 HTML 字符串
        title: 文件名前缀

    Returns:
        导出文件的相对路径
    """
    from app.services.browser_pool import managed_page

    safe_title = re.sub(r'[^\w\u4e00-\u9fff-]', '_', title)[:50]
    filename = f"{safe_title}_{uuid.uuid4().hex[:8]}.pdf"
    filepath = EXPORT_DIR / filename

    # 将 HTML 注入 print-pdf 模式的样式
    print_html = _inject_print_pdf_css(full_html)

    # 检测 HTML 结构类型
    is_deck_slide = 'class="deck-slide"' in full_html
    is_deck_page = 'class="deck-page"' in full_html
    is_webdeck = is_deck_slide or is_deck_page

    async with managed_page() as page:
        # 设置视口为标准幻灯片尺寸
        await page.set_viewport_size({"width": 1280, "height": 720})

        # 加载 HTML 内容
        await page.set_content(print_html, wait_until="networkidle")

        # 等待初始化
        await page.wait_for_timeout(2000)

        if is_webdeck:
            # webdeck 结构 - 支持 deck-slide 或 deck-page 两种 class
            slide_selector = '.deck-slide' if is_deck_slide else '.deck-page'
            slide_count = await page.evaluate(
                f"() => document.querySelectorAll('{slide_selector}').length"
            )
            logger.info(f"[Export] PDF 渲染: {slide_count} 页幻灯片 (webdeck: {slide_selector})")

            # 通过 JavaScript 修改 CSS
            await page.evaluate(f"""
                () => {{
                    // 设置 body 和 html 为 relative 定位，允许子元素垂直排列
                    document.body.style.position = 'relative';
                    document.body.style.height = 'auto';
                    document.body.style.overflow = 'visible';
                    document.documentElement.style.position = 'relative';
                    document.documentElement.style.height = 'auto';
                    document.documentElement.style.overflow = 'visible';

                    const container = document.getElementById('slides-container');
                    if (container) {{
                        container.style.position = 'relative';
                        container.style.height = 'auto';
                    }}

                    document.querySelectorAll('{slide_selector}').forEach((slide) => {{
                        slide.style.display = 'block';
                        slide.style.position = 'relative';
                        slide.style.top = 'auto';
                        slide.style.left = 'auto';
                        slide.style.transform = 'none';
                        slide.style.minHeight = '720px';
                        slide.style.height = 'auto';
                        slide.style.pageBreakAfter = 'always';
                    }});

                    const nav = document.getElementById('deck-nav-overlay');
                    const progress = document.getElementById('deck-progress');
                    if (nav) nav.style.display = 'none';
                    if (progress) progress.style.display = 'none';
                }}
            """)
        else:
            # reveal.js 结构
            slide_count = await page.evaluate("() => Reveal.getTotalSlides()")
            logger.info(f"[Export] PDF 渲染: {slide_count} 页幻灯片 (reveal.js)")

            # 通过 JavaScript 修改 CSS
            await page.evaluate("""
                () => {
                    const slidesContainer = document.querySelector('.reveal .slides');
                    if (slidesContainer) slidesContainer.style.position = 'relative';

                    document.querySelectorAll('.reveal .slides > section').forEach((slide) => {
                        slide.style.position = 'relative';
                        slide.style.display = 'block';
                        slide.style.opacity = '1';
                        slide.style.visibility = 'visible';
                        slide.style.top = 'auto';
                        slide.style.left = 'auto';
                        slide.style.transform = 'none';
                        slide.style.minHeight = '720px';
                        slide.style.pageBreakAfter = 'always';
                    });

                    const controls = document.querySelector('.reveal .controls');
                    const progress = document.querySelector('.reveal .progress');
                    const slideNumber = document.querySelector('.reveal .slide-number');
                    if (controls) controls.style.display = 'none';
                    if (progress) progress.style.display = 'none';
                    if (slideNumber) slideNumber.style.display = 'none';
                }
            """)

        # 等待布局稳定
        await page.wait_for_timeout(500)

        # 使用 page.pdf() 导出 PDF
        await page.pdf(
            path=str(filepath),
            width="1280px",
            height="720px",
            print_background=True,
            margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
        )

    logger.info(f"[Export] PDF 导出完成: {filepath}, {slide_count} 页")
    return f"exports/{filename}"


async def export_pptx_faithful(
    full_html: str,
    title: str,
    slide_count: int,
) -> str:
    """
    PPTX 保真导出 — 逐页截图后嵌入 python-pptx 生成图片型 PPTX。
    视觉效果与 reveal.js 渲染完全一致，但内容不可编辑。

    Args:
        full_html: 完整的 reveal.js HTML 字符串
        title: 文件名前缀
        slide_count: 幻灯片总页数

    Returns:
        导出文件的相对路径
    """
    from app.services.browser_pool import managed_page

    safe_title = re.sub(r'[^\w\u4e00-\u9fff-]', '_', title)[:50]
    filename = f"{safe_title}_{uuid.uuid4().hex[:8]}_faithful.pptx"
    filepath = EXPORT_DIR / filename

    # 截图每一页
    screenshots = []

    async with managed_page() as page:
        await page.set_viewport_size({"width": 1280, "height": 720})
        await page.set_content(full_html, wait_until="networkidle")
        await page.wait_for_timeout(2000)

        actual_count = await page.evaluate("() => Reveal.getTotalSlides()")
        total = min(slide_count, actual_count)

        for i in range(total):
            # 跳转到第 i 页
            await page.evaluate(f"() => Reveal.slide({i})")
            # 等待过渡动画完成（降低到100ms，reveal.js过渡动画通常<300ms）
            await page.wait_for_timeout(100)

            # 截图
            shot_path = EXPORT_DIR / f"_tmp_slide_{i}_{uuid.uuid4().hex[:6]}.png"
            await page.screenshot(path=str(shot_path), full_page=False)
            screenshots.append(shot_path)
            logger.debug(f"[Export] 截图第 {i+1} 页")

    # 使用 python-pptx 组装 PPTX
    prs = PptxPresentation()
    prs.slide_width = Emu(12192000)   # 1280 * 9525
    prs.slide_height = Emu(6858000)   # 720 * 9525

    blank_layout = prs.slide_layouts[6]  # 空白版式

    for i, shot_path in enumerate(screenshots):
        slide = prs.slides.add_slide(blank_layout)

        # 将截图铺满整个幻灯片
        slide.shapes.add_picture(
            str(shot_path),
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(str(filepath))

    # 清理临时截图
    for shot_path in screenshots:
        try:
            shot_path.unlink()
        except OSError:
            pass

    logger.info(f"[Export] PPTX 保真导出完成: {filepath}, {len(screenshots)} 页")
    return f"exports/{filename}"


async def export_pptx_editable(
    slides_data: list[dict[str, Any]],
    title: str,
    theme_id: str = "tech_dark",
) -> str:
    """
    PPTX 可编辑导出 — 解析 HTML 内容生成 python-pptx 文本型幻灯片。
    内容可编辑，但视觉效果与 reveal.js 有差异。

    Args:
        slides_data: [{index, html, speaker_notes}]
        title: 文件名前缀
        theme_id: 主题 ID

    Returns:
        导出文件的相对路径
    """
    safe_title = re.sub(r'[^\w\u4e00-\u9fff-]', '_', title)[:50]
    filename = f"{safe_title}_{uuid.uuid4().hex[:8]}_editable.pptx"
    filepath = EXPORT_DIR / filename

    theme = get_theme(theme_id)
    css = theme["css"]

    # 解析主题颜色
    bg_color = _hex_to_rgb(css.get("backgroundColor", "#ffffff"))
    heading_color = _hex_to_rgb(css.get("headingColor", "#1e40af"))
    text_color = _hex_to_rgb(css.get("color", "#1e293b"))

    prs = PptxPresentation()
    prs.slide_width = Emu(12192000)   # 1280 * 9525 EMU
    prs.slide_height = Emu(6858000)   # 720 * 9525 EMU

    blank_layout = prs.slide_layouts[6]  # 空白版式

    for slide_data in slides_data:
        html = slide_data.get("html", "")
        notes = slide_data.get("speaker_notes", "")

        pptx_slide = prs.slides.add_slide(blank_layout)

        # 设置背景色
        background = pptx_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg_color)

        # 解析 HTML 内容
        elements = _parse_slide_html(html)

        # 布局: 标题区 + 内容区
        y_cursor = Emu(457200)  # 0.48 inch 顶部边距
        left_margin = Emu(685800)  # 0.72 inch 左边距
        content_width = Emu(10820400)  # ~11.36 inch 内容宽度

        for elem in elements:
            tag = elem["tag"]
            text = elem["text"]

            if not text.strip():
                continue

            if tag in ("h1", "h2"):
                # 标题
                txBox = pptx_slide.shapes.add_textbox(
                    left_margin, y_cursor, content_width, Emu(685800)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(36 if tag == "h1" else 28)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*heading_color)
                y_cursor += Emu(800000)

            elif tag == "h3":
                txBox = pptx_slide.shapes.add_textbox(
                    left_margin, y_cursor, content_width, Emu(457200)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.color.rgb = RGBColor(*heading_color)
                y_cursor += Emu(571500)

            elif tag == "li":
                txBox = pptx_slide.shapes.add_textbox(
                    left_margin + Emu(228600), y_cursor, content_width - Emu(228600), Emu(342900)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"• {text}"
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(*text_color)
                y_cursor += Emu(400000)

            elif tag == "blockquote":
                txBox = pptx_slide.shapes.add_textbox(
                    left_margin + Emu(228600), y_cursor, content_width - Emu(457200), Emu(571500)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"「{text}」"
                p.font.size = Pt(16)
                p.font.italic = True
                p.font.color.rgb = RGBColor(*text_color)
                y_cursor += Emu(685800)

            else:
                # 普通段落
                txBox = pptx_slide.shapes.add_textbox(
                    left_margin, y_cursor, content_width, Emu(342900)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(*text_color)
                y_cursor += Emu(400000)

        # 添加演讲者备注
        if notes:
            notes_slide = pptx_slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    prs.save(str(filepath))
    logger.info(f"[Export] PPTX 可编辑导出完成: {filepath}, {len(slides_data)} 页")

    return f"exports/{filename}"


async def export_pptx_native(
    session: AsyncSession,
    presentation_id: str,
    slides_data: list[dict[str, Any]],
    title: str,
    theme_id: str = "tech_dark",
    user_id: str = DEFAULT_RUNTIME_USER_ID,
    workflow_package_id: str = OFFICIAL_NATIVE_ORCHESTRATOR,
) -> str:
    """通过官方 native-pptx-orchestrator workflow 导出 Native PPTX。"""
    result = await orchestrate_native_pptx_workflow(
        session,
        presentation_id=presentation_id,
        slides_data=slides_data,
        title=title,
        theme_id=theme_id,
        user_id=user_id,
        workflow_package_id=workflow_package_id,
        persist_artifact=True,
    )
    return result["file_path"]


async def orchestrate_native_pptx_workflow(
    session: AsyncSession,
    presentation_id: str,
    slides_data: list[dict[str, Any]],
    title: str,
    theme_id: str = "tech_dark",
    user_id: str = DEFAULT_RUNTIME_USER_ID,
    workflow_package_id: str = OFFICIAL_NATIVE_ORCHESTRATOR,
    persist_artifact: bool = True,
) -> dict[str, Any]:
    """Run a native PPTX workflow package and optionally persist the artifact."""
    safe_title = re.sub(r'[^\w\u4e00-\u9fff-]', '_', title)[:50]
    artifact_token = uuid.uuid4().hex[:8]
    filename = f"{safe_title}_{artifact_token}_native.pptx"
    preview_filename = f"{safe_title}_{artifact_token}_preview.html"
    filepath = EXPORT_DIR / filename
    preview_filepath = EXPORT_DIR / preview_filename

    runtime_result = await invoke_pptx_workflow_package(
        session,
        user_id,
        package_id=workflow_package_id,
        presentation_id=presentation_id,
        slides_data=slides_data,
        title=title,
        theme_id=theme_id,
    )

    if persist_artifact:
        filepath.write_bytes(runtime_result["pptx_content"])
        preview_filepath.write_text(runtime_result["html_preview_content"], encoding="utf-8")

    artifact_variant = await record_artifact_variant(
        session,
        user_id,
        package_id=runtime_result["workflow"]["package_id"],
        package_version=runtime_result["workflow"]["package_version"],
        variant_type="pptx-native",
        file_url=f"exports/{filename}" if persist_artifact else None,
        presentation_id=presentation_id,
        installed_plugin_id=runtime_result["workflow"].get("installed_plugin_id"),
        execution_log_id=runtime_result["workflow"].get("execution_log_id"),
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        metadata={
            "download_url": f"/static/exports/{filename}" if persist_artifact else None,
            "renderer_execution_log_id": runtime_result["renderer"].get("execution_log_id"),
        },
    )

    preview_variant = await record_artifact_variant(
        session,
        user_id,
        package_id=runtime_result["workflow"]["package_id"],
        package_version=runtime_result["workflow"]["package_version"],
        variant_type="html-preview",
        file_url=f"exports/{preview_filename}" if persist_artifact else None,
        presentation_id=presentation_id,
        installed_plugin_id=runtime_result["workflow"].get("installed_plugin_id"),
        execution_log_id=runtime_result["workflow"].get("execution_log_id"),
        mime_type="text/html",
        metadata={
            "download_url": f"/static/exports/{preview_filename}" if persist_artifact else None,
            "preview_renderer_execution_log_id": runtime_result["preview"].get("execution_log_id"),
            "artifact_role": "secondary_preview",
        },
    )
    await session.commit()

    logger.info(
        "[Export] PPTX Native 导出完成: %s, slides=%s, warnings=%s, workflow=%s@%s",
        filepath,
        runtime_result["renderer"].get("slideCount"),
        len(runtime_result["renderer"].get("warnings", [])),
        runtime_result["workflow"]["package_id"],
        runtime_result["workflow"]["package_version"],
    )
    return {
        "file_path": f"exports/{filename}" if persist_artifact else None,
        "download_url": f"/static/exports/{filename}" if persist_artifact else None,
        "preview_file_path": f"exports/{preview_filename}" if persist_artifact else None,
        "preview_download_url": f"/static/exports/{preview_filename}" if persist_artifact else None,
        "deck_spec": runtime_result["deck_spec"],
        "renderer": runtime_result["renderer"],
        "preview": runtime_result["preview"],
        "workflow": runtime_result["workflow"],
        "artifact_variant_id": artifact_variant.id,
        "html_artifact_variant_id": preview_variant.id,
    }


# ═══════════════ 内部辅助函数 ═══════════════


def _inject_print_pdf_css(html: str) -> str:
    """
    为 HTML 注入 print-pdf 所需的额外样式和配置。
    支持三种结构：
    1. reveal.js: .reveal .slides > section
    2. webdeck: .deck-slide + #slides-container
    3. webdeck: .deck-page (standalone section)

    让所有幻灯片在打印模式下垂直排列。
    """
    # 检测 HTML 结构类型
    has_deck_slide = 'class="deck-slide"' in html
    has_deck_page = 'class="deck-page"' in html
    has_reveal_slides = '.reveal .slides' in html or 'class="reveal"' in html

    print_css_parts = []

    if has_deck_slide or has_deck_page:
        # webdeck 结构 (artifact_composer.py 生成的 HTML)
        slide_class = '.deck-slide' if has_deck_slide else '.deck-page'
        print_css_parts.append(f"""
    <style>
      /* webdeck 结构打印修复 */
      #slides-container {{
        position: relative !important;
      }}
      {slide_class} {{
        display: block !important;
        position: relative !important;
        top: auto !important;
        left: auto !important;
        transform: none !important;
        page-break-after: always;
        min-height: 720px;
      }}
      /* 隐藏控制元素 */
      #deck-nav-overlay, #deck-progress {{
        display: none !important;
      }}
    </style>
        """)
    elif has_reveal_slides:
        # reveal.js 结构
        print_css_parts.append("""
    <style>
      /* 关键：让幻灯片容器使用相对定位，打破 absolute 堆叠 */
      .reveal .slides {
        position: relative !important;
      }
      /* 关键：让每张幻灯片使用相对定位，垂直流式排列 */
      .reveal .slides > section {
        position: relative !important;
        display: block !important;
        opacity: 1 !important;
        visibility: visible !important;
        page-break-after: always;
        min-height: 720px;
      }
      /* 隐藏控制元素 */
      .reveal .controls { display: none !important; }
      .reveal .progress { display: none !important; }
      .reveal .slide-number { display: none !important; }
    </style>
        """)

    if print_css_parts:
        print_css = "\n".join(print_css_parts)
        html = html.replace("</head>", f"{print_css}\n</head>")

    return html


def _parse_slide_html(html: str) -> list[dict[str, str]]:
    """
    简易 HTML 解析器 — 提取幻灯片中的结构化元素。
    将 HTML 转换为 [{tag, text}] 列表。

    支持标签: h1, h2, h3, p, li, blockquote, code
    """
    elements = []

    # 移除 <section> 标签包裹
    html = re.sub(r'</?section[^>]*>', '', html)

    # 按块级标签拆分
    patterns = [
        (r'<h1[^>]*>(.*?)</h1>', 'h1'),
        (r'<h2[^>]*>(.*?)</h2>', 'h2'),
        (r'<h3[^>]*>(.*?)</h3>', 'h3'),
        (r'<li[^>]*>(.*?)</li>', 'li'),
        (r'<blockquote[^>]*>(.*?)</blockquote>', 'blockquote'),
        (r'<p[^>]*>(.*?)</p>', 'p'),
    ]

    # 先标记所有匹配位置
    matches = []
    for pattern, tag in patterns:
        for m in re.finditer(pattern, html, re.DOTALL | re.IGNORECASE):
            text = re.sub(r'<[^>]+>', '', m.group(1)).strip()
            if text:
                matches.append((m.start(), tag, text))

    # 按出现位置排序
    matches.sort(key=lambda x: x[0])

    for _, tag, text in matches:
        elements.append({"tag": tag, "text": text})

    # 如果没有解析到任何结构化元素，整体提取纯文本
    if not elements:
        text = re.sub(r'<[^>]+>', '', html).strip()
        if text:
            elements.append({"tag": "p", "text": text})

    return elements


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """将 #rrggbb 颜色转为 (r, g, b) 整数元组。"""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) != 6:
        return (0, 0, 0)
    return (
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )
