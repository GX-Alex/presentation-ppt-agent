"""PPTX round-trip import utilities."""

from __future__ import annotations

import base64
from datetime import datetime
from pathlib import Path
import uuid
from typing import Any

from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.schemas.deck_spec import DeckSpec
from app.services.deckspec_preview_service import render_slide_sections_from_deckspec
from app.services.file_service import DEFAULT_USER_ID, resolve_file_reference
from app.services.ppt_service import create_presentation, persist_canonical_deckspec, save_slides
from app.services.theme_manager import get_theme
from app.services.user_settings_service import ensure_user
from app.models.tables import Asset, Task

PX_PER_INCH = 96
EMU_PER_INCH = 914400

IMAGE_MIME_BY_EXT = {
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
    "gif": "image/gif",
    "webp": "image/webp",
    "svg": "image/svg+xml",
}


class _ImportNodeFactory:
    def __init__(self, slide_index: int) -> None:
        self.slide_index = slide_index
        self.counter = 0

    def next_id(self, role: str) -> str:
        self.counter += 1
        safe_role = "-".join(part for part in role.lower().replace("_", "-").split("-") if part) or "node"
        return f"rt-slide-{self.slide_index + 1}-{safe_role}-{self.counter}"


async def import_pptx_as_presentation(
    session: AsyncSession,
    *,
    asset_id: str | None = None,
    file_url: str | None = None,
    title: str | None = None,
    theme_id: str | None = None,
    task_id: str | None = None,
    presentation_id: str | None = None,
    user_id: str = DEFAULT_USER_ID,
) -> dict[str, Any]:
    if not asset_id and not file_url:
        raise ValueError("必须提供 asset_id 或 file_url")

    await ensure_user(session, user_id=user_id)

    asset = None
    source_file_url = file_url
    if asset_id:
        result = await session.execute(select(Asset).where(Asset.id == asset_id))
        asset = result.scalar_one_or_none()
        if asset is None:
            raise ValueError("指定的资产不存在")
        if asset.file_type != "ppt" and not str(asset.file_url or "").lower().endswith(".pptx"):
            raise ValueError("指定资产不是 PPTX 文件")
        source_file_url = asset.file_url

    if not source_file_url:
        raise ValueError("未找到可导入的 PPTX 文件地址")

    file_path = resolve_file_reference(source_file_url)
    if not Path(file_path).exists():
        raise ValueError("PPTX 文件不存在或不可访问")

    import_title = title or (asset.title if asset is not None else Path(file_path).stem)
    task = await _get_or_create_import_task(session, task_id=task_id or (asset.task_id if asset else None), user_id=user_id, title=import_title)
    new_presentation_id = presentation_id or str(uuid.uuid4())

    deck_spec = import_pptx_to_deckspec(
        file_path,
        deck_id=new_presentation_id,
        title=import_title,
        theme_id=theme_id,
    )
    presentation_payload = build_presentation_payload_from_deckspec(deck_spec)

    await create_presentation(
        session,
        task_id=task.id,
        presentation_id=new_presentation_id,
        title=presentation_payload["title"],
        theme_id=presentation_payload["theme_id"],
        outline=presentation_payload["outline"],
        source_docs={
            "import_source": {
                "asset_id": asset.id if asset is not None else None,
                "file_url": source_file_url,
                "imported_at": datetime.utcnow().isoformat(),
                "kind": "pptx-roundtrip",
            }
        },
    )
    await save_slides(
        session,
        new_presentation_id,
        presentation_payload["slides"],
        refresh_canonical=False,
    )
    await persist_canonical_deckspec(
        session,
        new_presentation_id,
        deck_spec,
        source="pptx_roundtrip_importer",
        metadata={
            "asset_id": asset.id if asset is not None else None,
            "file_url": source_file_url,
            "task_id": task.id,
        },
    )

    return {
        "presentation_id": new_presentation_id,
        "task_id": task.id,
        "title": presentation_payload["title"],
        "theme_id": presentation_payload["theme_id"],
        "slide_count": len(deck_spec.slides),
        "outline": presentation_payload["outline"],
        "source": {
            "asset_id": asset.id if asset is not None else None,
            "file_url": source_file_url,
        },
        "deck_summary": {
            "deck_id": deck_spec.deck_id,
            "artifact_mode": deck_spec.artifact_mode,
            "slide_count": len(deck_spec.slides),
        },
    }


def import_pptx_to_deckspec(
    file_path: str,
    *,
    deck_id: str | None = None,
    title: str | None = None,
    theme_id: str | None = None,
) -> DeckSpec:
    prs = PptxPresentation(file_path)
    slide_payloads: list[dict[str, Any]] = []
    background_hex: str | None = None

    for slide_index, slide in enumerate(prs.slides):
        factory = _ImportNodeFactory(slide_index)
        nodes = _parse_slide_shapes(slide, factory)
        notes = _extract_notes(slide)
        slide_background = _extract_slide_background(slide)
        if background_hex is None and slide_background:
            background_hex = slide_background
        title_text = _identify_title(nodes, slide_index)
        page_type = _infer_page_type(nodes)
        slide_payloads.append(
            {
                "slide_id": f"imported-slide-{slide_index + 1}",
                "title": title_text,
                "page_type": page_type,
                "layout_id": f"roundtrip.{page_type}",
                "notes": notes,
                "nodes": nodes,
                "metadata": {
                    "source": "pptx-roundtrip-importer",
                    "slide_number": slide_index + 1,
                    "background": slide_background,
                },
            }
        )

    chosen_theme_id = theme_id or _guess_theme_id(background_hex)
    theme = get_theme(chosen_theme_id)
    theme_css = theme["css"]
    title_color = _find_first_text_color(slide_payloads, role="title")
    body_color = _find_first_text_color(slide_payloads, role="body")
    deck_title = title or _derive_deck_title(slide_payloads, Path(file_path).stem)
    slide_width = int(round(_emu_to_px(prs.slide_width)))
    slide_height = int(round(_emu_to_px(prs.slide_height)))

    deck_payload = {
        "deck_id": deck_id or str(uuid.uuid4()),
        "schema_version": "1.0.0",
        "revision": 1,
        "artifact_mode": "dual_render",
        "title": deck_title,
        "theme": {
            "theme_id": chosen_theme_id,
            "palette": {
                "background": background_hex or theme_css.get("backgroundColor", "#FFFFFF"),
                "foreground": body_color or theme_css.get("color", "#1E293B"),
                "accent": title_color or theme_css.get("headingColor", "#2563EB"),
                "muted": theme_css.get("linkColor", theme_css.get("accentColor", "#64748B")),
            },
            "typography": {
                "heading_font": _strip_font_family(theme_css.get("headingFontFamily")) or "Aptos Display",
                "body_font": _strip_font_family(theme_css.get("fontFamily")) or "Aptos",
                "mono_font": _strip_font_family(theme_css.get("codeFontFamily")) or "Cascadia Code",
            },
            "spacing": {
                "base_unit": 8,
                "section_gap": 24,
                "item_gap": 12,
            },
            "custom": {
                "source": "pptx-roundtrip-importer",
                "source_file": Path(file_path).name,
            },
        },
        "slide_size": {
            "width": slide_width,
            "height": slide_height,
            "unit": "px",
        },
        "slides": slide_payloads,
        "metadata": {
            "source": "pptx-roundtrip-importer",
            "source_file_path": file_path,
            "imported_slide_count": len(slide_payloads),
        },
    }
    return DeckSpec.model_validate(deck_payload)


def build_presentation_payload_from_deckspec(deck_spec: DeckSpec) -> dict[str, Any]:
    sections = render_slide_sections_from_deckspec(deck_spec)
    slides_data = [
        {
            "index": index,
            "type": slide.page_type,
            "html": sections[index],
            "speaker_notes": slide.notes,
        }
        for index, slide in enumerate(deck_spec.slides)
    ]
    outline = [
        {
            "index": index,
            "title": slide.title,
            "type": slide.page_type,
            "bullets": _outline_bullets(slide),
            "speaker_notes": slide.notes,
        }
        for index, slide in enumerate(deck_spec.slides)
    ]
    return {
        "title": deck_spec.title,
        "theme_id": deck_spec.theme.theme_id,
        "outline": outline,
        "slides": slides_data,
    }


async def _get_or_create_import_task(
    session: AsyncSession,
    *,
    task_id: str | None,
    user_id: str,
    title: str,
) -> Task:
    if task_id:
        result = await session.execute(select(Task).where(Task.id == task_id))
        task = result.scalar_one_or_none()
        if task is not None:
            return task

    task = Task(
        id=str(uuid.uuid4()),
        user_id=user_id,
        title=title,
        status="active",
        intent="ppt",
    )
    session.add(task)
    await session.commit()
    return task


def _parse_slide_shapes(slide: Any, factory: _ImportNodeFactory) -> list[dict[str, Any]]:
    nodes: list[dict[str, Any]] = []
    for shape in slide.shapes:
        nodes.extend(_parse_shape(shape, factory, 0, 0))
    nodes.sort(key=lambda item: (item["bbox"]["y"], item["bbox"]["x"]))
    return nodes


def _parse_shape(shape: Any, factory: _ImportNodeFactory, offset_x: int, offset_y: int) -> list[dict[str, Any]]:
    shape_type = getattr(shape, "shape_type", None)
    bbox = _shape_bbox(shape, offset_x, offset_y)
    nodes: list[dict[str, Any]] = []

    if shape_type == MSO_SHAPE_TYPE.GROUP:
        child_nodes: list[dict[str, Any]] = []
        for child in getattr(shape, "shapes", []):
            child_nodes.extend(_parse_shape(child, factory, offset_x + int(shape.left), offset_y + int(shape.top)))
        return child_nodes

    if getattr(shape, "has_table", False):
        rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
        col_count = max((len(row) for row in rows), default=1)
        nodes.append(
            {
                "node_id": factory.next_id("table"),
                "kind": "table",
                "role": "table",
                "bbox": bbox,
                "content": {
                    "rows": rows,
                    "headerRows": 1 if rows else 0,
                    "colWidths": [bbox["w"] / max(col_count, 1) for _ in range(max(col_count, 1))],
                },
                "style": {
                    "fontSize": 12,
                    "lineColor": _color_from_line(shape),
                    "fillColor": _color_from_fill(shape),
                },
                "children": [],
            }
        )
        return nodes

    if getattr(shape, "has_chart", False):
        chart_payload = _parse_chart(shape)
        if chart_payload:
            nodes.append(
                {
                    "node_id": factory.next_id("chart"),
                    "kind": "chart",
                    "role": "chart",
                    "bbox": bbox,
                    "content": chart_payload,
                    "style": {
                        "color": _color_from_line(shape),
                        "fillColor": _color_from_fill(shape),
                    },
                    "children": [],
                }
            )
        return nodes

    image_node = _parse_image(shape, factory, bbox)
    if image_node is not None:
        return [image_node]

    text_nodes = _parse_text_shape(shape, factory, bbox)
    if text_nodes:
        return text_nodes

    shape_node = _parse_generic_shape(shape, factory, bbox)
    return [shape_node] if shape_node is not None else []


def _parse_text_shape(shape: Any, factory: _ImportNodeFactory, bbox: dict[str, float]) -> list[dict[str, Any]]:
    if not getattr(shape, "has_text_frame", False):
        return []
    paragraphs = getattr(shape.text_frame, "paragraphs", [])
    text = "\n".join(paragraph.text.strip() for paragraph in paragraphs if paragraph.text.strip())
    if not text:
        return []

    runs: list[dict[str, Any]] = []
    font_sizes: list[float] = []
    is_bold = False
    is_italic = False
    text_color: str | None = None
    alignment: str | None = None
    for paragraph_index, paragraph in enumerate(paragraphs):
        para_text = paragraph.text.strip()
        if not para_text:
            continue
        alignment = alignment or _map_alignment(getattr(paragraph, "alignment", None))
        paragraph_runs = list(getattr(paragraph, "runs", []) or [])
        if paragraph_runs:
            for run_index, run in enumerate(paragraph_runs):
                run_text = (run.text or "").strip("\n")
                if not run_text:
                    continue
                font = getattr(run, "font", None)
                run_color = _color_from_font(font)
                text_color = text_color or run_color
                font_size = _font_size_from_font(font)
                if font_size:
                    font_sizes.append(font_size)
                is_bold = is_bold or bool(getattr(font, "bold", False))
                is_italic = is_italic or bool(getattr(font, "italic", False))
                runs.append(
                    {
                        "text": run_text,
                        "bold": bool(getattr(font, "bold", False)),
                        "italic": bool(getattr(font, "italic", False)),
                        "color": run_color,
                        "bullet": _is_bulleted_paragraph(paragraph),
                        "indentLevel": getattr(paragraph, "level", 0) or 0,
                        "breakLine": False,
                    }
                )
                if run_index == len(paragraph_runs) - 1 and paragraph_index < len(paragraphs) - 1:
                    runs[-1]["breakLine"] = True
        else:
            runs.append(
                {
                    "text": para_text,
                    "bold": False,
                    "italic": False,
                    "color": None,
                    "bullet": _is_bulleted_paragraph(paragraph),
                    "indentLevel": getattr(paragraph, "level", 0) or 0,
                    "breakLine": paragraph_index < len(paragraphs) - 1,
                }
            )

    shape_style = _extract_shape_style(shape)
    nodes: list[dict[str, Any]] = []
    if _has_visible_box(shape_style):
        nodes.append(
            {
                "node_id": factory.next_id("textbox-bg"),
                "kind": "shape",
                "role": "textbox.background",
                "bbox": bbox,
                "content": {"shapeType": "roundedRect"},
                "style": shape_style,
                "children": [],
            }
        )

    nodes.append(
        {
            "node_id": factory.next_id("text"),
            "kind": "text",
            "role": "body",
            "bbox": bbox,
            "content": {"text": text, "runs": runs},
            "style": {
                "fontSize": max(font_sizes) if font_sizes else _infer_font_size_from_bbox(bbox),
                "bold": is_bold,
                "italic": is_italic,
                "color": text_color,
                "align": alignment or "left",
            },
            "children": [],
        }
    )
    return nodes


def _parse_image(shape: Any, factory: _ImportNodeFactory, bbox: dict[str, float]) -> dict[str, Any] | None:
    shape_type = getattr(shape, "shape_type", None)
    if shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None
    image = getattr(shape, "image", None)
    if image is None:
        return None
    ext = str(getattr(image, "ext", "png") or "png").lower()
    mime_type = IMAGE_MIME_BY_EXT.get(ext, "image/png")
    blob = getattr(image, "blob", None)
    if not blob:
        return None
    encoded = base64.b64encode(blob).decode("ascii")
    return {
        "node_id": factory.next_id("image"),
        "kind": "image",
        "role": "image",
        "bbox": bbox,
        "content": {
            "data": f"data:{mime_type};base64,{encoded}",
            "alt": getattr(shape, "name", "image"),
        },
        "style": {},
        "children": [],
    }


def _parse_generic_shape(shape: Any, factory: _ImportNodeFactory, bbox: dict[str, float]) -> dict[str, Any] | None:
    style = _extract_shape_style(shape)
    if not _has_visible_box(style):
        return None
    return {
        "node_id": factory.next_id("shape"),
        "kind": "shape",
        "role": "shape",
        "bbox": bbox,
        "content": {"shapeType": "roundedRect"},
        "style": style,
        "children": [],
    }


def _parse_chart(shape: Any) -> dict[str, Any] | None:
    chart = getattr(shape, "chart", None)
    if chart is None:
        return None

    categories: list[str] = []
    try:
        categories = [str(category.label) for category in chart.plots[0].categories]
    except Exception:
        categories = []

    series_payload = []
    try:
        for series in chart.series:
            values = [float(value or 0) for value in list(series.values)]
            series_payload.append({"name": str(series.name or "Series"), "values": values})
    except Exception:
        series_payload = []

    if not categories or not series_payload:
        return None
    return {
        "chartType": _map_chart_type(getattr(chart, "chart_type", None)),
        "categories": categories,
        "series": series_payload,
        "showLegend": len(series_payload) > 1,
        "showValue": True,
    }


def _identify_title(nodes: list[dict[str, Any]], slide_index: int) -> str:
    text_nodes = [node for node in nodes if node["kind"] == "text" and str(node.get("content", {}).get("text") or "").strip()]
    if not text_nodes:
        return f"Slide {slide_index + 1}"

    def score(node: dict[str, Any]) -> tuple[float, float]:
        font_size = float(node.get("style", {}).get("fontSize") or 16)
        top = float(node.get("bbox", {}).get("y") or 0)
        return (font_size, -top)

    title_node = max(text_nodes, key=score)
    title_node["role"] = "title"
    title_node.setdefault("style", {})
    title_node["style"]["bold"] = True
    title_node["style"]["fontSize"] = max(float(title_node["style"].get("fontSize") or 16), 24)

    title_text = str(title_node.get("content", {}).get("text") or "").strip().splitlines()[0]
    for node in text_nodes:
        node.setdefault("style", {})
        node["style"].setdefault("align", "left")
        if node is not title_node:
            node["role"] = node.get("role") or "body"
    return title_text or f"Slide {slide_index + 1}"


def _infer_page_type(nodes: list[dict[str, Any]]) -> str:
    kinds = {node["kind"] for node in nodes}
    if "table" in kinds:
        return "table"
    if "chart" in kinds:
        return "chart"
    image_count = sum(1 for node in nodes if node["kind"] == "image")
    text_nodes = [node for node in nodes if node["kind"] == "text" and node.get("role") != "title"]
    if image_count and text_nodes:
        return "image-text"
    x_positions = sorted(float(node.get("bbox", {}).get("x") or 0) for node in text_nodes)
    if len(x_positions) >= 2 and (max(x_positions) - min(x_positions)) > 220:
        return "two-column"
    return "content"


def _outline_bullets(slide: Any) -> list[str]:
    bullets: list[str] = []
    for node in slide.nodes:
        if node.kind != "text" or node.role == "title":
            continue
        text = str(node.content.get("text") or "").strip()
        if not text:
            continue
        bullets.append(text.splitlines()[0][:80])
        if len(bullets) >= 3:
            break
    return bullets


def _derive_deck_title(slides: list[dict[str, Any]], fallback: str) -> str:
    if slides:
        title = str(slides[0].get("title") or "").strip()
        if title:
            return title
    return fallback or "Imported Presentation"


def _extract_notes(slide: Any) -> str:
    try:
        text = slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""
    if text.lower().startswith("click to add notes"):
        return ""
    return text


def _extract_slide_background(slide: Any) -> str | None:
    try:
        return _color_from_fore_color(slide.background.fill.fore_color)
    except Exception:
        return None


def _find_first_text_color(slides: list[dict[str, Any]], *, role: str) -> str | None:
    for slide in slides:
        for node in slide.get("nodes", []):
            if node.get("kind") == "text" and node.get("role") == role:
                color = node.get("style", {}).get("color")
                if color:
                    return str(color)
    return None


def _guess_theme_id(background_hex: str | None) -> str:
    if not background_hex:
        return "business_light"
    rgb = background_hex.lstrip("#")
    if len(rgb) != 6:
        return "business_light"
    red = int(rgb[0:2], 16)
    green = int(rgb[2:4], 16)
    blue = int(rgb[4:6], 16)
    luminance = (0.2126 * red + 0.7152 * green + 0.0722 * blue) / 255
    return "tech_dark" if luminance < 0.45 else "business_light"


def _shape_bbox(shape: Any, offset_x: int, offset_y: int) -> dict[str, float]:
    left = offset_x + int(getattr(shape, "left", 0) or 0)
    top = offset_y + int(getattr(shape, "top", 0) or 0)
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return {
        "x": round(_emu_to_px(left), 2),
        "y": round(_emu_to_px(top), 2),
        "w": max(round(_emu_to_px(width), 2), 1.0),
        "h": max(round(_emu_to_px(height), 2), 1.0),
    }


def _emu_to_px(value: int | float) -> float:
    return float(value) / EMU_PER_INCH * PX_PER_INCH


def _font_size_from_font(font: Any) -> float | None:
    size = getattr(font, "size", None)
    try:
        return float(size.pt) if size is not None else None
    except Exception:
        return None


def _infer_font_size_from_bbox(bbox: dict[str, float]) -> float:
    return max(14.0, min(28.0, bbox["h"] * 0.32))


def _map_alignment(alignment: Any) -> str | None:
    if alignment is None:
        return None
    value = str(alignment).lower()
    if "center" in value:
        return "center"
    if "right" in value:
        return "right"
    if "justify" in value:
        return "justify"
    return "left"


def _is_bulleted_paragraph(paragraph: Any) -> bool:
    text = str(getattr(paragraph, "text", "") or "").strip()
    if getattr(paragraph, "level", 0):
        return True
    return text.startswith(("•", "-", "▪", "◦", "·"))


def _extract_shape_style(shape: Any) -> dict[str, Any]:
    fill_color = _color_from_fill(shape)
    line_color = _color_from_line(shape)
    line_width = _line_width(shape)
    transparency = _fill_transparency(shape)
    return {
        "fillColor": fill_color,
        "lineColor": line_color,
        "lineWidth": line_width,
        "transparency": transparency,
    }


def _has_visible_box(style: dict[str, Any]) -> bool:
    return bool(style.get("fillColor") or style.get("lineColor"))


def _color_from_fill(shape: Any) -> str | None:
    try:
        return _color_from_fore_color(shape.fill.fore_color)
    except Exception:
        return None


def _color_from_line(shape: Any) -> str | None:
    try:
        return _color_from_fore_color(shape.line.color)
    except Exception:
        return None


def _line_width(shape: Any) -> float | None:
    try:
        width = getattr(shape.line, "width", None)
        return float(width.pt) if width is not None else None
    except Exception:
        return None


def _fill_transparency(shape: Any) -> float | None:
    try:
        transparency = getattr(shape.fill, "transparency", None)
        return float(transparency) if transparency is not None else None
    except Exception:
        return None


def _color_from_font(font: Any) -> str | None:
    try:
        return _color_from_fore_color(font.color)
    except Exception:
        return None


def _color_from_fore_color(fore_color: Any) -> str | None:
    rgb = getattr(fore_color, "rgb", None)
    if rgb is None:
        return None
    return f"#{str(rgb)}"


def _map_chart_type(chart_type: Any) -> str:
    value = str(chart_type).lower()
    if "line" in value:
        return "line"
    if "pie" in value:
        return "pie"
    return "bar"


def _strip_font_family(font_family: str | None) -> str | None:
    if not font_family:
        return None
    first = font_family.split(",")[0].strip().strip("'\"")
    return first or None